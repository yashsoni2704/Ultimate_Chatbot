"""
Document loader — detects file type and routes to the correct extractor.

Supported formats
-----------------
PDF   → Docling  (primary, AI-based layout + OCR)  + PyPDFLoader fallback
DOCX  → Docling
PPTX  → Docling
XLSX  → Docling (each table row → own Document for precise retrieval)
XLS   → xlrd  (legacy, Docling does not read .xls)
CSV   → CSVLoader
TXT   → plain read
RTF   → striprtf

Why Docling instead of Tesseract
---------------------------------
Docling uses IBM's layout-understanding models (DocLayNet + TableFormer) with
an integrated EasyOCR/RapidOCR engine. It reconstructs reading order, merges
multi-column text correctly, and handles tables as structured data — all the
things Tesseract broke into random character fragments.

Docling is lazy-imported on first use so startup time is unaffected.
"""

# ── stdlib ───────────────────────────────────────────────────────────────────
import io
import os
import zipfile
from datetime import datetime

# ── LangChain ────────────────────────────────────────────────────────────────
from langchain_core.documents import Document
from langchain_community.document_loaders import CSVLoader
from langchain_text_splitters.character import RecursiveCharacterTextSplitter

# ── project ──────────────────────────────────────────────────────────────────
from config import Config
from utils.logger import get_logger
from utils.chunker import chunk_documents, _enrich
from utils.embeddings import EmbeddingManager

logger = get_logger(__name__)

# ── LangSmith ────────────────────────────────────────────────────────────────
if Config.LANGCHAIN_TRACING_V2:
    os.environ["LANGCHAIN_TRACING_V2"] = "true"
    os.environ["LANGCHAIN_API_KEY"]     = Config.LANGCHAIN_API_KEY
    os.environ["LANGCHAIN_PROJECT"]     = Config.LANGCHAIN_PROJECT

try:
    from langsmith import traceable as _traceable
    _LANGSMITH_AVAILABLE = True
except ImportError:
    def _traceable(**_kwargs):
        def _decorator(fn):
            return fn
        return _decorator
    _LANGSMITH_AVAILABLE = False


# ════════════════════════════════════════════════════════════════════════════
# Docling singleton — lazy-loaded on first use
# ════════════════════════════════════════════════════════════════════════════

_docling_converter = None   # DocumentConverter instance (expensive to create)


def _get_docling_converter():
    """
    Return a shared DocumentConverter, creating it on the first call.

    PipelineOptions are configured for maximum quality:
      • do_ocr=True          — always run OCR on image regions / scanned pages
      • do_table_structure=True — reconstruct table cells (not just raw text)
    """
    global _docling_converter
    if _docling_converter is None:
        try:
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.pipeline.standard_pdf_pipeline import PdfPipelineOptions

            pdf_opts = PdfPipelineOptions()
            pdf_opts.do_ocr             = True   # OCR any image-only regions
            pdf_opts.do_table_structure = True   # TableFormer → clean table text

            _docling_converter = DocumentConverter(
                format_options={"pdf": PdfFormatOption(pipeline_options=pdf_opts)}
            )
            logger.info("✅ Docling DocumentConverter initialised (OCR + TableFormer enabled)")
        except Exception as exc:
            logger.error(f"❌ Docling failed to initialise: {exc}")
            raise RuntimeError(
                f"Docling could not be initialised: {exc}. "
                "Run:  pip install docling"
            ) from exc
    return _docling_converter


def _docling_convert_to_documents(file_path: str, file_type: str) -> list:
    """
    Convert any Docling-supported file to a list of LangChain Documents.

    Strategy
    --------
    • Convert with Docling → get a DoclingDocument object
    • Export each page (or the whole file for non-PDF) to Markdown
    • Wrap every non-empty page as its own Document with metadata

    For PDFs we iterate page-by-page so the page number is preserved in
    metadata, matching what pdfplumber did.  For DOCX/PPTX/XLSX the whole
    file is one Markdown export split on Docling's page breaks.
    """
    converter = _get_docling_converter()

    logger.info(f"  Docling converting: {os.path.basename(file_path)}")
    result = converter.convert(file_path)
    doc    = result.document       # DoclingDocument

    documents = []

    if file_type == "pdf":
        # Export each page individually to preserve page-level metadata
        total_pages = len(doc.pages) if hasattr(doc, "pages") and doc.pages else 1
        for page_no in range(total_pages):
            try:
                page_md = doc.export_to_markdown(page_no=page_no).strip()
            except TypeError:
                # Older Docling versions don't support page_no — fall back to full
                page_md = doc.export_to_markdown().strip()
            if page_md:
                documents.append(Document(
                    page_content=page_md,
                    metadata={
                        "source":    file_path,
                        "page":      page_no,
                        "file_type": "pdf",
                        "extractor": "docling",
                    },
                ))
        # If per-page export failed, fall back to full-document export
        if not documents:
            full_md = doc.export_to_markdown().strip()
            if full_md:
                documents.append(Document(
                    page_content=full_md,
                    metadata={
                        "source":    file_path,
                        "file_type": "pdf",
                        "extractor": "docling",
                    },
                ))
    else:
        # DOCX / PPTX — one Markdown export, split on page-break markers
        full_md = doc.export_to_markdown().strip()
        if not full_md:
            return []

        # Docling uses "<!-- page break -->" to separate pages/slides
        pages = [p.strip() for p in full_md.split("<!-- page break -->") if p.strip()]
        if not pages:
            pages = [full_md]

        for i, page_md in enumerate(pages):
            documents.append(Document(
                page_content=page_md,
                metadata={
                    "source":    file_path,
                    "page":      i,
                    "file_type": file_type,
                    "extractor": "docling",
                },
            ))

    logger.info(
        f"  Docling extracted {len(documents)} section(s) "
        f"from '{os.path.basename(file_path)}'"
    )
    return documents


# ════════════════════════════════════════════════════════════════════════════
# PDF loader  — Docling primary + PyPDFLoader fallback
# ════════════════════════════════════════════════════════════════════════════

def _load_pdf(file_path: str) -> list:
    """
    Extract text from a PDF.

    Pipeline
    --------
    1. Try Docling (AI layout + OCR + table reconstruction).
       Handles scanned PDFs, multi-column layouts, embedded images, tables.
    2. If Docling fails for any reason, fall back to PyPDFLoader
       so ingestion never silently stops.
    """
    try:
        documents = _docling_convert_to_documents(file_path, "pdf")
        if documents:
            logger.info(
                f"✅ Docling extracted {len(documents)} page(s) "
                f"from '{os.path.basename(file_path)}'"
            )
            return documents
        logger.warning("  Docling returned no content — falling back to PyPDFLoader")
    except Exception as exc:
        logger.warning(f"  Docling failed ({exc}) — falling back to PyPDFLoader")

    # Fallback
    from langchain_community.document_loaders import PyPDFLoader
    docs = PyPDFLoader(file_path).load()
    for d in docs:
        d.metadata["extractor"] = "pypdf_fallback"
    logger.info(
        f"✅ PyPDFLoader (fallback) extracted {len(docs)} page(s) "
        f"from '{os.path.basename(file_path)}'"
    )
    return docs


# ════════════════════════════════════════════════════════════════════════════
# DOCX loader  — Docling primary + python-docx fallback
# ════════════════════════════════════════════════════════════════════════════

def _load_docx(file_path: str) -> list:
    """
    Load a Word document via Docling (handles embedded images, tables,
    headers, footers, text boxes).  Falls back to python-docx for plain text.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".doc":
        raise Exception(
            "Legacy .doc format is not supported. "
            "Please save as .docx and re-upload."
        )

    try:
        documents = _docling_convert_to_documents(file_path, "docx")
        if documents:
            logger.info(
                f"✅ Docling extracted {len(documents)} section(s) "
                f"from '{os.path.basename(file_path)}'"
            )
            return documents
        logger.warning("  Docling returned no content — falling back to python-docx")
    except Exception as exc:
        logger.warning(f"  Docling failed ({exc}) — falling back to python-docx")

    # Fallback: python-docx
    try:
        import docx as _docx
    except ImportError:
        raise ImportError("python-docx is required: pip install python-docx")

    try:
        doc = _docx.Document(file_path)
    except Exception as e:
        raise Exception(f"Could not open Word document: {e}")

    text_parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            text_parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                text_parts.append(" | ".join(cells))

    full_text = "\n".join(text_parts).strip()
    if not full_text:
        return []
    return [Document(
        page_content=full_text,
        metadata={"source": file_path, "file_type": "docx", "extractor": "python-docx-fallback"},
    )]


# ════════════════════════════════════════════════════════════════════════════
# PPTX loader  — Docling primary + python-pptx fallback
# ════════════════════════════════════════════════════════════════════════════

def _load_pptx_text(file_path: str) -> list:
    """
    Load a PowerPoint file via Docling (handles slide images, diagrams,
    speaker notes, table cells).  Falls back to python-pptx for text shapes.
    """
    try:
        documents = _docling_convert_to_documents(file_path, "pptx")
        if documents:
            logger.info(
                f"✅ Docling extracted {len(documents)} slide(s) "
                f"from '{os.path.basename(file_path)}'"
            )
            return documents
        logger.warning("  Docling returned no content — falling back to python-pptx")
    except Exception as exc:
        logger.warning(f"  Docling failed ({exc}) — falling back to python-pptx")

    # Fallback: python-pptx (text shapes only, no images)
    from pptx import Presentation
    prs       = Presentation(file_path)
    documents = []
    for i, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    text_parts.append(t)
        content = "\n".join(text_parts).strip()
        if content:
            documents.append(Document(
                page_content=content,
                metadata={
                    "source":            file_path,
                    "slide_number":      i,
                    "file_type":         "pptx",
                    "extractor":         "python-pptx-fallback",
                },
            ))
    logger.info(f"[PPTX fallback] {len(documents)}/{len(prs.slides)} slides")
    return documents


# kept for any external callers
def load_pptx_with_ocr(file_path: str) -> list:
    return _load_pptx_text(file_path)


# ════════════════════════════════════════════════════════════════════════════
# XLSX loader  — Docling primary + openpyxl fallback (one row per Document)
# ════════════════════════════════════════════════════════════════════════════

def _load_xlsx(file_path: str) -> list:
    """
    Load an Excel file.

    Strategy
    --------
    Docling primary  → exports the whole workbook as Markdown, then we split
    on Docling's table/page-break markers so each sheet becomes its own
    Document.  This preserves column headers and cell relationships in a
    way that Tesseract + raw cell scanning never could.

    openpyxl fallback → one Document per data row (old behaviour, reliable
    for purely tabular files without images).

    Legacy .xls is always handled by xlrd (Docling doesn't read .xls).
    """
    ext = os.path.splitext(file_path)[1].lower()

    # ── Legacy .xls via xlrd ─────────────────────────────────────────────────
    if ext == ".xls":
        return _load_xls_xlrd(file_path)

    # ── .xlsx via Docling ────────────────────────────────────────────────────
    try:
        documents = _docling_convert_to_documents(file_path, "xlsx")
        if documents:
            logger.info(
                f"✅ Docling extracted {len(documents)} section(s) "
                f"from '{os.path.basename(file_path)}'"
            )
            return documents
        logger.warning("  Docling returned no XLSX content — falling back to openpyxl")
    except Exception as exc:
        logger.warning(f"  Docling failed on XLSX ({exc}) — falling back to openpyxl")

    # ── openpyxl fallback — one Document per data row ────────────────────────
    return _load_xlsx_openpyxl(file_path)


def _load_xlsx_openpyxl(file_path: str) -> list:
    """One Document per data row via openpyxl (fallback)."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required: pip install openpyxl")

    wb        = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    documents = []

    for sheet_name in wb.sheetnames:
        ws       = wb[sheet_name]
        all_rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            all_rows.append(cells)

        if not all_rows:
            continue

        # First non-empty row → headers
        headers, header_row_idx = [], 0
        for idx, row_cells in enumerate(all_rows):
            non_empty = [c for c in row_cells if c]
            if non_empty:
                headers       = [c if c else f"Column{i+1}" for i, c in enumerate(row_cells)]
                header_row_idx = idx
                break
        if not headers:
            continue

        for row_idx, row_cells in enumerate(all_rows[header_row_idx + 1:],
                                            start=header_row_idx + 1):
            if not any(c for c in row_cells):
                continue
            pairs    = [f"{headers[ci] if ci < len(headers) else f'Column{ci+1}'}: {v or '-'}"
                        for ci, v in enumerate(row_cells)]
            row_text = " | ".join(pairs)
            documents.append(Document(
                page_content=row_text,
                metadata={
                    "source":     file_path,
                    "sheet":      sheet_name,
                    "row_number": row_idx,
                    "file_type":  "xlsx",
                    "extractor":  "openpyxl-fallback",
                },
            ))

    wb.close()
    return documents


def _load_xls_xlrd(file_path: str) -> list:
    """Legacy .xls via xlrd — one Document per data row."""
    try:
        import xlrd
    except ImportError:
        raise ImportError("xlrd is required: pip install xlrd")

    workbook  = xlrd.open_workbook(file_path)
    documents = []

    for sheet_name in workbook.sheet_names():
        sheet = workbook.sheet_by_name(sheet_name)
        if sheet.nrows == 0:
            continue

        headers, header_row_idx = [], 0
        for row_idx in range(sheet.nrows):
            candidate = [str(sheet.cell_value(row_idx, col)).strip()
                         for col in range(sheet.ncols)]
            candidate = [c for c in candidate if c]
            if candidate:
                headers       = [str(sheet.cell_value(row_idx, col)).strip() or f"Column{col+1}"
                                 for col in range(sheet.ncols)]
                header_row_idx = row_idx
                break
        if not headers:
            continue

        for row_idx in range(header_row_idx + 1, sheet.nrows):
            raw = [str(sheet.cell_value(row_idx, col)).strip() for col in range(sheet.ncols)]
            if not any(raw):
                continue
            pairs    = [f"{headers[ci] if ci < len(headers) else f'Column{ci+1}'}: {v or '-'}"
                        for ci, v in enumerate(raw)]
            row_text = " | ".join(pairs)
            documents.append(Document(
                page_content=row_text,
                metadata={
                    "source":     file_path,
                    "sheet":      sheet_name,
                    "row_number": row_idx,
                    "file_type":  "xls",
                },
            ))

    return documents


# ════════════════════════════════════════════════════════════════════════════
# Simple format loaders (CSV / TXT / RTF) — unchanged, no OCR needed
# ════════════════════════════════════════════════════════════════════════════

def _load_txt(file_path: str) -> list:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read().strip()
            if content:
                return [Document(
                    page_content=content,
                    metadata={"source": file_path, "file_type": "txt"},
                )]
            return []
        except UnicodeDecodeError:
            continue
    raise ValueError(f"Could not decode text file: {file_path}")


def _load_rtf(file_path: str) -> list:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError("striprtf is required: pip install striprtf")

    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                raw = f.read()
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError(f"Could not decode RTF file: {file_path}")

    content = rtf_to_text(raw).strip()
    if not content:
        return []
    return [Document(
        page_content=content,
        metadata={"source": file_path, "file_type": "rtf"},
    )]


# ════════════════════════════════════════════════════════════════════════════
# MAIN ROUTER
# ════════════════════════════════════════════════════════════════════════════

def load_document(file_path: str) -> list:
    """
    Detect file type by extension and route to the correct loader.
    Returns a list of LangChain Document objects.
    """
    extension = os.path.splitext(file_path)[1].lower()
    logger.info(
        f"Detected file extension: '{extension}' "
        f"for '{os.path.basename(file_path)}'"
    )

    if extension == ".pdf":
        logger.info(f"Loading PDF: {os.path.basename(file_path)}")
        documents = _load_pdf(file_path)

    elif extension == ".csv":
        loader    = CSVLoader(file_path)
        documents = loader.load()

    elif extension in (".ppt", ".pptx"):
        logger.info(f"Loading PowerPoint: {os.path.basename(file_path)}")
        documents = _load_pptx_text(file_path)

    elif extension in (".docx", ".doc"):
        logger.info(f"Loading Word document: {os.path.basename(file_path)}")
        documents = _load_docx(file_path)

    elif extension in (".xlsx", ".xls"):
        logger.info(f"Loading Excel: {os.path.basename(file_path)}")
        documents = _load_xlsx(file_path)

    elif extension == ".txt":
        logger.info(f"Loading plain text: {os.path.basename(file_path)}")
        documents = _load_txt(file_path)

    elif extension == ".rtf":
        logger.info(f"Loading RTF: {os.path.basename(file_path)}")
        documents = _load_rtf(file_path)

    else:
        raise Exception(
            f"Unsupported file type '{extension}'. "
            "Supported: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, CSV, TXT, RTF"
        )

    logger.info(f"Loaded {len(documents)} document section(s) from {file_path}")
    return documents


# ════════════════════════════════════════════════════════════════════════════
# PROCESS PIPELINE
# ════════════════════════════════════════════════════════════════════════════

def create_chunks(documents):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=Config.CHUNK_SIZE,
        chunk_overlap=Config.CHUNK_OVERLAP,
    )
    return splitter.split_documents(documents)


@_traceable(
    name="process_document",
    run_type="chain",
    tags=["ingestion", "docmind"],
    metadata={"component": "loader"},
)
def process_document(file_path: str) -> str:
    """
    Full ingestion pipeline:
      1. Detect type → load via Docling (or lightweight loader for CSV/TXT/RTF)
      2. Chunk + contextual enrichment
      3. Embed → blue/green vector store
    """
    filename  = os.path.basename(file_path)
    extension = os.path.splitext(file_path)[1].lower().lstrip(".")

    logger.info("=" * 80)
    logger.info("📄 DOCUMENT PROCESSING STARTED")
    logger.info(f"  File      : {file_path}")
    logger.info(f"  Timestamp : {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    try:
        # Step 1 — Load
        logger.info("Step 1: Loading document...")
        documents = _traced_load(file_path, filename=filename, file_type=extension)
        logger.info(f"✅ Loaded {len(documents)} document section(s)")

        if not documents:
            raise Exception("No readable content found.")

        for i, doc in enumerate(documents):
            logger.info(f"  Section {i+1}: {len(doc.page_content)} chars")

        # Step 2 — Chunk + enrich
        # Excel: each Document is already one row — skip text splitter,
        # just apply contextual enrichment so row identity is preserved.
        if extension in ("xlsx", "xls"):
            logger.info("Step 2: Excel — enriching rows as whole units (no splitting)...")
            chunks = [_enrich(doc) for doc in documents]
            logger.info(f"✅ Created {len(chunks)} enriched row chunks")
        else:
            logger.info(
                f"Step 2: Chunking + enrichment "
                f"(size={Config.CHUNK_SIZE}, overlap={Config.CHUNK_OVERLAP})..."
            )
            chunks = _traced_chunk(documents, filename=filename)
            logger.info(f"✅ Created {len(chunks)} enriched chunks")

        if not chunks:
            raise Exception(
                f"No text could be extracted from '{filename}'. "
                "The file may be empty or use an unsupported encoding."
            )

        for i, chunk in enumerate(chunks[:3]):
            logger.info(f"  Chunk {i+1} preview: {chunk.page_content[:120]!r}")

        # Step 3 — Embed + save via blue/green pipeline
        logger.info("Step 3: Embedding into vector store (blue/green pipeline)...")
        embedding_manager = EmbeddingManager()
        result = embedding_manager.create_vector_store(chunks, source_path=file_path)

        if result is None:
            return f"'{filename}' is already in the knowledge base. No duplicate added."

        logger.info("✅ Vector store updated successfully")
        logger.info("📄 DOCUMENT PROCESSING COMPLETED")
        logger.info("=" * 80)

        return f"Document processed successfully. {len(chunks)} chunks created."

    except Exception as e:
        logger.error(f"❌ Error processing document: {str(e)}")
        logger.error("=" * 80)
        raise


@_traceable(name="load_document_file", run_type="tool", tags=["ingestion", "load"])
def _traced_load(file_path: str, filename: str = "", file_type: str = "") -> list:
    return load_document(file_path)


@_traceable(name="chunk_documents", run_type="tool", tags=["ingestion", "chunking"])
def _traced_chunk(documents: list, filename: str = "") -> list:
    return chunk_documents(documents)


# ════════════════════════════════════════════════════════════════════════════
# URL INGESTION PIPELINE  — unchanged logic, Docling handles downloaded files
# ════════════════════════════════════════════════════════════════════════════

def load_url(url: str) -> str:
    """
    Smart ingestion pipeline for any URL.

    Classification → route:
      GOOGLE_DRIVE  → gdown download      → process_document()
      ONEDRIVE      → requests download   → process_document()
      DIRECT_FILE   → requests download   → process_document()
      WEB_PAGE      → Playwright scrape   → chunk + embed directly

    Downloaded files go through process_document() which now uses Docling,
    so scanned PDFs from Drive/OneDrive also benefit from AI OCR.
    """
    from utils.url_resolver import classify, UrlType
    from utils.scraper import scrape_url
    from utils.embeddings import (
        delete_document_chunks,
        _load_registry,
        _save_registry,
    )

    url = url.strip()

    logger.info("=" * 80)
    logger.info("🌐  URL INGESTION STARTED")
    logger.info(f"  URL       : {url}")
    logger.info(f"  Timestamp : {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    classification = classify(url)
    logger.info(f"  Classified as: {classification}")

    try:
        # De-duplicate: remove previous ingestion of the same URL
        registry = [r for r in _load_registry() if r is not None]
        existing = [r for r in registry if r.get("filename") == url]
        if existing:
            logger.info("  Already ingested — removing old data for re-ingestion...")
            deleted = delete_document_chunks(url)
            logger.info(f"  Deleted {deleted} old vector(s)")
            registry = [r for r in registry if r.get("filename") != url]
            _save_registry(registry)

        if classification.url_type == UrlType.WEB_PAGE:
            return _ingest_web_page(url, scrape_url, chunk_documents, EmbeddingManager)
        else:
            return _ingest_file_from_url(url, classification)

    except Exception as e:
        logger.error(f"  ❌ Error ingesting URL: {str(e)}")
        logger.error("=" * 80)
        raise


def _ingest_web_page(url, scrape_url_fn, chunk_fn, EmbeddingManagerCls) -> str:
    logger.info("  Step 1: Scraping web page with Playwright...")
    documents = scrape_url_fn(url)
    logger.info(f"  ✅ Scraped {len(documents)} section(s)")

    if not documents:
        raise Exception("No content extracted from URL.")

    for doc in documents:
        if doc.metadata is None:
            doc.metadata = {"source": url, "source_type": "url"}

    logger.info("  Step 2: Chunking + contextual enrichment...")
    chunks = chunk_fn(documents)
    logger.info(f"  ✅ Created {len(chunks)} enriched chunks")

    if not chunks:
        raise Exception("No text content could be extracted from the URL.")

    logger.info("  Step 3: Embedding into vector store...")
    embedding_manager = EmbeddingManagerCls()
    result = embedding_manager.create_vector_store(chunks, source_path=url)

    if result is None:
        return f"URL already in knowledge base: {url}"

    logger.info("  ✅ Vectors stored")
    logger.info("🌐  URL INGESTION COMPLETE")
    logger.info("=" * 80)
    return f"Web page scraped and ingested successfully. {len(chunks)} chunks created."


def _ingest_file_from_url(url: str, classification) -> str:
    from utils.url_resolver import UrlType
    from utils.downloader import (
        download_google_drive,
        download_onedrive,
        download_direct_file,
    )
    from utils.embeddings import _load_registry, _save_registry

    local_path = None

    try:
        logger.info(f"  Step 1: Downloading {classification.label}...")

        if classification.url_type == UrlType.GOOGLE_DRIVE:
            local_path = download_google_drive(url, extension_hint=classification.extension_hint)
        elif classification.url_type == UrlType.ONEDRIVE:
            local_path = download_onedrive(url, extension_hint=classification.extension_hint)
        elif classification.url_type == UrlType.DIRECT_FILE:
            local_path = download_direct_file(url, extension_hint=classification.extension_hint)
        else:
            raise RuntimeError(f"Unexpected URL type: {classification.url_type}")

        logger.info(f"  ✅ Downloaded to: {local_path}")

        logger.info(f"  Step 2/3: Processing via Docling pipeline...")
        process_document(local_path)

        logger.info(f"  Step 4: Re-keying registry source to original URL...")
        _rekey_registry_source(local_path, url)

        logger.info("🌐  URL INGESTION COMPLETE")
        logger.info("=" * 80)
        return (
            f"{classification.label} downloaded and ingested successfully. "
            f"Source recorded as: {url}"
        )

    finally:
        if local_path and os.path.exists(local_path):
            try:
                os.remove(local_path)
                logger.info(f"  🧹 Cleaned up temp file: {local_path}")
            except OSError as rm_err:
                logger.warning(f"  ⚠️  Could not remove temp file {local_path}: {rm_err}")


def _rekey_registry_source(old_source: str, new_source: str) -> None:
    from utils.embeddings import _load_registry, _save_registry
    try:
        registry = [r for r in _load_registry() if r is not None]
        updated  = False
        for entry in registry:
            if entry.get("filename") == old_source:
                entry["filename"] = new_source
                updated = True
            if entry.get("source") == old_source:
                entry["source"] = new_source
                updated = True
        if updated:
            _save_registry(registry)
            logger.info(f"  ✅ Registry re-keyed: {os.path.basename(old_source)} → {new_source}")
        else:
            logger.warning(
                f"  ⚠️  Re-key: no registry entry for '{old_source}'. "
                "KB table may show local path instead of URL."
            )
    except Exception as e:
        logger.warning(f"  ⚠️  Registry re-key failed (non-fatal): {e}")
